# -*- coding: utf-8 -*-
"""Gaziantep Üniversitesi Rektörlüğü — resmi dilde kapsamlı ürün sunumu."""
from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from tempfile import NamedTemporaryFile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_pptx import (
    CYAN,
    INK,
    LINE,
    MUTED,
    NAVY,
    NAVY2,
    NAVY3,
    SOFT,
    TEAL,
    WHITE,
    add_multiline,
    add_shot,
    add_textbox,
    bullet_block,
    footer_bar,
    header_bar,
    matrix_2x2,
    process_row,
    rounded_rect,
    section_title,
    stack_layers,
)

HERE = Path(__file__).resolve().parent
SHOTS = HERE / "shots"
OUT = HERE / "KampusteyimAPP_GAUN_Rektorluk_Sunumu.pptx"
OUT_ALT = HERE / "KampusteyimAPP_GAUN_Urun_Odakli_Sunum.pptx"
ROOT = HERE.parents[1]

# Dinamik sayfa sayacı
_PAGE = {"n": 0, "total": 18}


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slide_layouts[6]


def dark_background(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.4), Inches(-1.2), Inches(5.4), Inches(5.4)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = NAVY3
    glow.line.fill.background()


def frame(slide, prs, kicker: str, title: str):
    _PAGE["n"] += 1
    header_bar(slide, prs, "Gaziantep Üniversitesi Rektörlüğü")
    section_title(slide, kicker, title)
    footer_bar(slide, prs, _PAGE["n"], _PAGE["total"])


def phone(slide, image_name, left, top, width, height):
    rounded_rect(slide, left, top, width, height, NAVY2)
    inner = Inches(0.12)
    path = SHOTS / image_name
    if path.exists():
        add_shot(
            slide,
            path,
            left + inner,
            top + inner,
            width - inner * 2,
            height - inner * 2,
        )


def two_col_bullets(slide, left_title, left_items, right_title, right_items):
    rounded_rect(slide, Inches(0.65), Inches(1.75), Inches(5.9), Inches(4.85), NAVY2)
    add_textbox(
        slide,
        Inches(0.95),
        Inches(2.05),
        Inches(5.3),
        Inches(0.4),
        left_title,
        size=18,
        bold=True,
        color=CYAN,
    )
    bullet_block(slide, Inches(0.95), Inches(2.6), Inches(5.2), left_items, size=13)

    rounded_rect(
        slide, Inches(6.8), Inches(1.75), Inches(5.9), Inches(4.85), SOFT, line=LINE
    )
    add_textbox(
        slide,
        Inches(7.1),
        Inches(2.05),
        Inches(5.3),
        Inches(0.4),
        right_title,
        size=18,
        bold=True,
        color=NAVY2,
    )
    bullet_block(slide, Inches(7.1), Inches(2.6), Inches(5.2), right_items, size=13)


def feature_grid(slide, cards: list[tuple[str, str]], cols=3):
    n = len(cards)
    rows = (n + cols - 1) // cols
    card_w = Inches(12.0 / cols - 0.15)
    card_h = Inches(4.7 / max(rows, 1) - 0.12)
    for i, (title, body) in enumerate(cards):
        col, row = i % cols, i // cols
        x = Inches(0.65 + col * (12.0 / cols))
        y = Inches(1.75 + row * (4.85 / max(rows, 1)))
        dark = (row + col) % 2 == 0
        fill = NAVY2 if dark else SOFT
        tc = CYAN if dark else NAVY2
        bc = RGBColor(0xD5, 0xE4, 0xEE) if dark else MUTED
        rounded_rect(
            slide, x, y, card_w, card_h, fill, line=None if dark else LINE
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            y + Inches(0.18),
            card_w - Inches(0.4),
            Inches(0.35),
            title,
            size=14,
            bold=True,
            color=tc,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            y + Inches(0.58),
            card_w - Inches(0.4),
            card_h - Inches(0.75),
            body,
            size=11,
            color=bc,
        )


def cover(prs, blank):
    slide = prs.slides.add_slide(blank)
    dark_background(slide, prs)
    logo = SHOTS / "logo_kampus.png"
    if logo.exists():
        slide.shapes.add_picture(
            str(logo), Inches(0.85), Inches(0.55), width=Inches(0.95)
        )
    add_textbox(
        slide,
        Inches(0.85),
        Inches(1.7),
        Inches(7.6),
        Inches(0.4),
        "Gaziantep Üniversitesi Rektörlüğü’ne sunulmak üzere",
        size=14,
        color=CYAN,
    )
    add_textbox(
        slide,
        Inches(0.85),
        Inches(2.25),
        Inches(7.8),
        Inches(0.55),
        "KampüsteyimAPP",
        size=40,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(0.85),
        Inches(2.95),
        Inches(7.8),
        Inches(1.15),
        "Üniversite kampüsünün doğrulanmış\ndijital iletişim ve yönetim platformu",
        size=26,
        bold=True,
        color=CYAN,
    )
    add_textbox(
        slide,
        Inches(0.85),
        Inches(4.4),
        Inches(7.5),
        Inches(1.1),
        "Öğrenci, resmi öğrenci topluluğu, firma ve kurumsal süreçleri "
        "tek güvenli uygulamada birleştiren; Gaziantep Üniversitesi’nde "
        "öncü pilot uygulama olarak konumlandırılması önerilen kampüs platformudur.",
        size=15,
        color=RGBColor(0xD5, 0xE4, 0xEE),
    )
    add_textbox(
        slide,
        Inches(0.85),
        Inches(6.35),
        Inches(7.5),
        Inches(0.35),
        "Hazırlayan: AYS Tech  ·  Kayra Çatalkaya  ·  info@kampusteyim.app",
        size=12,
        bold=True,
        color=CYAN,
    )
    phone(slide, "03_feed.png", Inches(9.15), Inches(0.7), Inches(3.2), Inches(6.15))


def purpose(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "01 · SUNUMUN AMACI",
        "Rektörlüğümüze arz edilen talep ve kapsam",
    )
    two_col_bullets(
        slide,
        "Arz edilen hususlar",
        [
            "KampüsteyimAPP’in GAÜN’de öncü / pilot dijital kampüs modeli olarak değerlendirilmesi",
            "Öğrenci toplulukları ile koordinasyonun Rektörlük ve ilgili birimler nezdinde desteklenmesi",
            "Doğrulanmış resmi iletişim kanalının üniversite markasına uygun biçimde konumlandırılması",
            "Sağlık Kültür ve Spor (SKS) süreçlerinin modernizasyonuna katkı zemininin açılması",
            "Yüz yüze veya çevrim içi brifing ile teknik–idari detayların netleştirilmesi",
        ],
        "Platformun kurumsal niteliği",
        [
            "Ürün, AYS Tech tarafından geliştirilmiş bir kampüs sosyal ağıdır",
            "Kayıt ve hesap süreçleri öğrenci doğrulamasına dayalıdır",
            "Topluluk ve firma hesapları yönetici onayı ile açılır",
            "İçerik, şikâyet ve kısıtlama süreçleri merkezi yönetim paneli ile denetlenir",
            "KVKK / gizlilik tercihleri ve hesap güvenliği ürün mimarisinin parçasıdır",
        ],
    )


def vision(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "02 · KURUMSAL VİZYON",
        "Tek kampüs, tek doğrulanmış dijital deneyim",
    )
    matrix_2x2(
        slide,
        Inches(0.65),
        Inches(1.8),
        Inches(7.3),
        Inches(4.75),
        [
            (
                "Öğrenci",
                "Kampüs gündemini takip eder; paylaşım, etkileşim, etkinlik başvurusu ve kariyer araçlarına erişir.",
            ),
            (
                "Resmi topluluk",
                "Doğrulanmış kimlikle duyuru, etkinlik, üye iletişimi ve görünürlüğü kurumsal biçimde yönetir.",
            ),
            (
                "Üniversite",
                "Parçalı mesajlaşma ve sosyal medya dağınıklığı yerine ölçülebilir, denetlenebilir bir kanal elde eder.",
            ),
            (
                "Ekosistem",
                "Sosyal yaşam, kültür–spor, kariyer ve işveren etkileşimi aynı güvenli çatı altında buluşur.",
            ),
        ],
    )
    rounded_rect(slide, Inches(8.2), Inches(1.8), Inches(4.5), Inches(4.75), SOFT, line=LINE)
    add_textbox(
        slide,
        Inches(8.5),
        Inches(2.15),
        Inches(3.95),
        Inches(0.7),
        "Parçalı iletişimin yerine",
        size=20,
        bold=True,
        color=NAVY2,
    )
    bullet_block(
        slide,
        Inches(8.5),
        Inches(3.05),
        Inches(3.9),
        [
            "Tek uygulama ve tek hesap",
            "Doğrulanmış kimlik katmanları",
            "Kampüse özel içerik ve filtreler",
            "Rol bazlı yetkilendirme",
            "Merkezî moderasyon ve raporlama",
            "Mobil erişim ve anlık bildirim",
        ],
        size=14,
    )


def actors(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(slide, prs, "03 · PAYDAŞLAR", "Platformdaki rol ve hesap türleri")
    feature_grid(
        slide,
        [
            (
                "Öğrenci hesabı",
                "E-posta doğrulama, öğrenci belgesi / kart ile kayıt; akış, hikâye, Reels, etkinlik, CV ve staj araçları.",
            ),
            (
                "Topluluk hesabı",
                "Yalnızca yönetici onayıyla açılır; logo zorunluluğu, altın/mavi rozet, duyuru ve etkinlik yetkisi.",
            ),
            (
                "Firma hesabı",
                "Kayıt kapalıdır; admin açar ve teslim eder. Web paneli ile ilan, CV tarama, mail ve teklif.",
            ),
            (
                "Kampüs elçisi",
                "Topluluğu temsil eden öğrencilere özel rozet ve görünürlük; aidiyet göstergesi.",
            ),
            (
                "Süper admin / personel",
                "Kullanıcı, içerik, şikâyet, yasak, rozet, firma/topluluk açılışı ve sistem ayarları.",
            ),
            (
                "İzleyici / misafir",
                "Kısıtlı görüntüleme; etkileşim ve başvuru için giriş zorunluluğu (AuthGate).",
            ),
        ],
    )


def registration(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "04 · KAYIT VE DOĞRULAMA",
        "Kampüse özel güvenli üyelik süreci",
    )
    process_row(
        slide,
        Inches(0.55),
        Inches(1.85),
        Inches(12.2),
        [
            ("E-posta kodu", "Doğrulama bileti"),
            ("Profil", "Ad · kampüs · kullanıcı adı"),
            ("Belge", "Kart / öğrenci belgesi"),
            ("Onay", "Bekleyen / onaylı hesap"),
        ],
    )
    rounded_rect(slide, Inches(0.65), Inches(3.45), Inches(5.9), Inches(3.15), NAVY2)
    add_textbox(
        slide,
        Inches(0.95),
        Inches(3.65),
        Inches(5.3),
        Inches(0.35),
        "Süreç ayrıntıları",
        size=16,
        bold=True,
        color=CYAN,
    )
    bullet_block(
        slide,
        Inches(0.95),
        Inches(4.15),
        Inches(5.2),
        [
            "Kayıt öncesi e-posta doğrulama kodu zorunludur",
            "Kullanıcı adı AI moderasyonu ve teklik kontrolü",
            "Öğrenci kimlik / belge ile hesap durumu yönetimi",
            "Bekleyen hesaplar sınırlı erişimle onaya alınır",
            "KVKK onayı ve pazarlama izni kayıt anında alınır",
        ],
        size=12,
    )
    rounded_rect(
        slide, Inches(6.8), Inches(3.45), Inches(5.9), Inches(3.15), SOFT, line=LINE
    )
    add_textbox(
        slide,
        Inches(7.1),
        Inches(3.65),
        Inches(5.3),
        Inches(0.35),
        "Kurumsal güvence",
        size=16,
        bold=True,
        color=NAVY2,
    )
    bullet_block(
        slide,
        Inches(7.1),
        Inches(4.15),
        Inches(5.2),
        [
            "Sahte veya yetkisiz hesapların azaltılması",
            "Üniversiteye özgü kimlik alanı (şehir / üniversite)",
            "Silinen hesapların yeniden canlandırılmaması",
            "Oturum bütünlüğü ve güvenli oturum yenileme",
            "Admin tarafından hesap silme / devre dışı bırakma",
        ],
        size=12,
    )


def feed(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "05 · KAMPÜS AKIŞI",
        "Gönderiler, etkileşim ve kapsam filtreleri",
    )
    phone(slide, "03_feed.png", Inches(0.65), Inches(1.7), Inches(3.15), Inches(5.0))
    add_textbox(
        slide,
        Inches(4.15),
        Inches(1.75),
        Inches(8.5),
        Inches(0.55),
        "Öğrenciler metin, görsel, video ve belge paylaşabilir; beğeni, yorum, "
        "yorum yanıtı, yorum sabitleme, repost ve # / @ sistemleri tam sosyal ağ "
        "standartında çalışır.",
        size=13,
        color=INK,
    )
    bullet_block(
        slide,
        Inches(4.15),
        Inches(2.55),
        Inches(8.4),
        [
            "Kapsam filtreleri: Tümü · Şehrim · Üniversitem (anlık snapshot ile saniyelik güncelleme)",
            "Yeni gönderiler denormalize şehir/üniversite alanlarıyla doğru filtreye düşer",
            "Hashtag keşfi ve popüler etiketler; bahsetme bildirimleri",
            "Çift dokunuşla beğeni; medya görüntüleyici ve güvenli ağ görselleri",
            "Gizli hesap / engelleme / aramadan gizlenme kurallarına uygun görünürlük",
            "Sponsorlu içerik ve reklam yerleşimleri (yönetilebilir)",
            "İçerik Guard (yerel güvenlik + bulut moderasyon) ile zararlı paylaşımın engellenmesi",
        ],
        size=13,
    )


def stories_reels(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "06 · HİKÂYE VE REELS",
        "Kısa ömürlü ve kısa video kampüs medyası",
    )
    feature_grid(
        slide,
        [
            (
                "Hikâye halkası",
                "Görülmemiş / görülmüş halka; canlı dinleme ile saniyelik yansıma; optimistic yayın.",
            ),
            (
                "Stream-first video",
                "Tam dosya indirme beklenmez; progressive oynatma ve sonraki öğe ısıtması.",
            ),
            (
                "Kampüs Reels",
                "Dikey kaydırma, N±1 pencere, izlenme kaydı; oturumda sabit sıralama.",
            ),
            (
                "Etkileşim",
                "Beğeni, yorum, takip isteği, izleyici listesi; sahip için görüntülenme.",
            ),
            (
                "Paylaşım",
                "Kamera / galeri; video postların Reels’e yansıması (ayna).",
            ),
            (
                "Kurumsal kullanım",
                "Topluluk ve firma marka anlatımı için kısa, denetlenebilir format.",
            ),
        ],
        cols=3,
    )


def discovery(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "07 · KEŞİF VE PROFİL",
        "Arama, öneriler ve kişisel vitrin",
    )
    phone(slide, "05_search.png", Inches(0.65), Inches(1.7), Inches(3.05), Inches(5.0))
    phone(slide, "07_profile.png", Inches(3.95), Inches(1.7), Inches(3.05), Inches(5.0))
    rounded_rect(slide, Inches(7.3), Inches(1.7), Inches(5.35), Inches(5.0), SOFT, line=LINE)
    add_textbox(
        slide,
        Inches(7.6),
        Inches(2.0),
        Inches(4.8),
        Inches(0.4),
        "Profil ve keşif ayrıntıları",
        size=17,
        bold=True,
        color=NAVY2,
    )
    bullet_block(
        slide,
        Inches(7.6),
        Inches(2.55),
        Inches(4.8),
        [
            "İsim, kullanıcı adı, e-posta ve öğrenci no ile arama",
            "Önerilen kişiler: kampüs affinity + ağırlıklı karışım",
            "Biyografi, bağlantı, foto, doğrulama rozetleri",
            "Gönderi / Reels ızgarası; kullanıcı adı değişince eski içerik taşınır",
            "Gizlilik: gizli hesap, takip isteği, engelleme, aramadan gizle",
            "İzleyici modu: paylaşım kapalı, salt görüntüleme",
            "Üniversite rozet etiketleri (ör. GAÜN, GİBTÜ)",
        ],
        size=12,
    )


def events(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "08 · ETKİNLİK YÖNETİMİ",
        "Duyurudan katılıma uçtan uca operasyon",
    )
    phone(slide, "06_events.png", Inches(0.65), Inches(1.7), Inches(3.2), Inches(5.0))
    process_row(
        slide,
        Inches(4.2),
        Inches(1.85),
        Inches(8.4),
        [
            ("Oluştur", "Tarih · yer · görsel"),
            ("Yayınla", "Akış · hikâye"),
            ("Başvuru", "Kota · son tarih"),
            ("Yönet", "Onay · doluluk"),
        ],
    )
    rounded_rect(slide, Inches(4.2), Inches(3.4), Inches(8.4), Inches(3.25), SOFT, line=LINE)
    add_textbox(
        slide,
        Inches(4.55),
        Inches(3.65),
        Inches(7.8),
        Inches(0.35),
        "Kurumsal ve topluluk kazanımları",
        size=16,
        bold=True,
        color=NAVY2,
    )
    bullet_block(
        slide,
        Inches(4.55),
        Inches(4.15),
        Inches(7.7),
        [
            "Dağınık form, WhatsApp ve e-posta trafiğinin azaltılması",
            "Kontenjan, son başvuru ve hedef kitle bilgisinin tek kartta toplanması",
            "Başvuru durumlarının canlı izlenmesi; süre dolunca otomatik kapanış desteği",
            "Öğrencinin etkinliği keşfetmesi, paylaşması ve hatırlatma alması",
            "App Links / Universal Links ile etkinlik bağlantısının mağazaya düşmesi",
        ],
        size=13,
    )


def communities(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "09 · RESMİ TOPLULUKLAR",
        "Doğrulanmış öğrenci topluluklarının dijital merkezi",
    )
    feature_grid(
        slide,
        [
            (
                "Resmi kimlik",
                "Admin onayı, logo zorunluluğu, altın/mavi rozet; sahte hesaplardan ayrışma.",
            ),
            (
                "İçerik yetkisi",
                "Duyuru, gönderi, hikâye, Reels ve etkinlik kartı yayını.",
            ),
            (
                "Üye ve takip",
                "Takip, gizli hesap, takip isteği; üye görünürlüğü ve iletişim.",
            ),
            (
                "Elçilik",
                "Kampüs elçisi rozeti; öğrenci–topluluk aidiyet bağları.",
            ),
            (
                "İş ortakları",
                "About bölümünde partner vitrini (logo, açıklama, boyut).",
            ),
            (
                "Kurumsal denetim",
                "Şikâyet, kısıtlama ve Guard ile içerik güvenliği.",
            ),
        ],
    )


def career(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "10 · KARİYER KATMANI",
        "CV-AI, Staj-AI ve firma paneli",
    )
    stack_layers(
        slide,
        Inches(0.7),
        Inches(1.85),
        Inches(12.0),
        [
            ("CV-AI", "Öğrenci ATS uyumlu özgeçmiş üretir; dışa aktarım ve güncelleme", TEAL),
            ("Staj-AI", "İlan listesi, başvuru, teklif kutusu; CV yoksa yönlendirme", NAVY2),
            ("Firma paneli (web)", "İlan CRUD, öğrenci tarama, mail, direkt teklif, başvuru sıralama AI", NAVY3),
            ("Eşleşme", "Yetkinlik ve profil sinyalleriyle nitelikli aday–işveren köprüsü", RGBColor(0x17, 0x8A, 0x84)),
            ("Kurumsal değer", "Mezun istihdamı ve staj süreçlerinde ölçülebilir dijital kanal", TEAL),
        ],
    )


def notifications_plus(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "11 · BİLDİRİM, PLUS VE MARKET",
        "Anlık iletişim ve sürdürülebilir görünürlük",
    )
    two_col_bullets(
        slide,
        "Bildirim sistemi",
        [
            "Uygulama içi bildirim merkezi",
            "FCM push: beğeni, yorum, takip, repost, ilan, teklif",
            "Bildirim tercihleri (kullanıcı kontrolü)",
            "Kısıtlama / yasak bildirimleri (push + mail + inbox)",
            "Sessiz / odaklanmış bildirim politikası",
        ],
        "Plus ve Market",
        [
            "KampüsteyimPlus: yeşil tick ve öncelikli deneyim",
            "Web Market: app.kampusteyim.app/market",
            "Çoklu ay planları ve ödeme entegrasyonu",
            "Öğrenci ve marka görünürlüğü için abonelik katmanı",
            "İndirme: kampusteyim.app/get.html (akıllı mağaza yönlendirme)",
        ],
    )


def admin_moderation(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "12 · YÖNETİM VE MODERASYON",
        "Merkezî denetim paneli",
    )
    feature_grid(
        slide,
        [
            (
                "Kullanıcı yönetimi",
                "Profil düzenleme, şifre sıfırlama, hesap silme, rol ve rozet atama.",
            ),
            (
                "Topluluk / firma",
                "Hesap açma–teslim, topluluk rozeti, firma paneli erişimi.",
            ),
            (
                "İçerik denetimi",
                "Tüm gönderi/yorum; silme, gizleme, hard delete.",
            ),
            (
                "Şikâyetler",
                "Gönderi, yorum, hesap şikâyetleri; AI ön inceleme desteği.",
            ),
            (
                "Kısıtlama",
                "Geçici/kalıcı paylaşım yasağı, ban; süre ve gerekçe ile bildirim.",
            ),
            (
                "Personel rolleri",
                "Staff rol kimlikleri ile yetki ayrımı; süper admin koruması.",
            ),
        ],
    )


def privacy_security(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "13 · GİZLİLİK VE GÜVENLİK",
        "Kurumsal kullanıma uygun teknik ve hukuki zemin",
    )
    matrix_2x2(
        slide,
        Inches(0.65),
        Inches(1.8),
        Inches(12.0),
        Inches(4.75),
        [
            (
                "Kişisel veri",
                "KVKK onayı, gizlilik politikası, hesap silme akışı; silinen hesapların yeniden canlandırılmaması.",
            ),
            (
                "Oturum güvenliği",
                "Kalıcı oturum, bütünlük kontrolü, sessiz yenileme; oturum karışmasının önlenmesi.",
            ),
            (
                "İçerik güvenliği",
                "Yerel güvenlik filtresi, bulut Guard, raporlama ve yaptırım zinciri.",
            ),
            (
                "Altyapı",
                "Firebase Auth / Firestore / Storage / Cloud Functions; App Links ve Universal Links.",
            ),
        ],
    )


def sks(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "14 · SKS VE KÜLTÜR–SPOR",
        "Dinamik, mobil ve ölçülebilir süreç önerisi",
    )
    two_col_bullets(
        slide,
        "KampüsteyimAPP katmanı",
        [
            "Topluluk ve etkinlik görünürlüğünün mobil erişimi",
            "Başvuru, kontenjan ve hatırlatma süreçleri",
            "Kültür–spor faaliyetlerinin öğrenci odaklı duyurusu",
            "Geri bildirim ve katılımın dijital izlenebilirliği",
            "Yetkilendirilmiş yönetim panelleri",
        ],
        "Kuruma sunulan fayda",
        [
            "Eski / dağınık altyapının modern mobil deneyime evrilmesi",
            "Öğrenciye güncel ve tek kanaldan erişim",
            "Ölçülebilir katılım ve operasyon verisi",
            "Üniversite markasına uygun resmi dijital yüzey",
            "Mevcut idari süreçlere saygılı, aşamalı entegrasyon",
        ],
    )


def roadmap(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "15 · PİLOT YOL HARİTASI",
        "GAÜN’de kontrollü yaygınlaştırma önerisi",
    )
    process_row(
        slide,
        Inches(0.55),
        Inches(1.95),
        Inches(12.2),
        [
            ("1. Kapsam", "Birimler ve KPI"),
            ("2. Doğrulama", "Topluluk / öğrenci"),
            ("3. Pilot", "Sınırlı kullanım"),
            ("4. Ölçüm", "Geri bildirim"),
            ("5. Yaygınlaştırma", "Üniversite geneli"),
        ],
    )
    rounded_rect(slide, Inches(0.8), Inches(3.55), Inches(11.7), Inches(2.7), NAVY2)
    add_textbox(
        slide,
        Inches(1.15),
        Inches(3.9),
        Inches(11.0),
        Inches(0.4),
        "Önerilen pilot göstergeleri",
        size=18,
        bold=True,
        color=CYAN,
    )
    add_multiline(
        slide,
        Inches(1.15),
        Inches(4.5),
        Inches(11.0),
        Inches(1.4),
        [
            "• Aktif öğrenci ve doğrulanmış topluluk hesabı sayısı",
            "• Etkinlik yayın / başvuru / doluluk oranları",
            "• Duyuru erişimi ve bildirim etkileşimi",
            "• Şikâyet çözüm süreleri ve moderasyon yükü",
            "• Kariyer / staj kanalı kullanım metrikleri (opsiyonel faz)",
        ],
        size=14,
        color=WHITE,
        gap_pt=4,
    )


def feature_inventory(prs, blank):
    slide = prs.slides.add_slide(blank)
    frame(
        slide,
        prs,
        "16 · ÖZELLİK ENVANTERİ",
        "Ürünün başlıca yeteneklerinin özeti",
    )
    cols = [
        (
            "Sosyal",
            [
                "Akış & filtreler",
                "Gönderi / medya / PDF",
                "Beğeni · yorum · repost",
                "Hashtag · mention",
                "Hikâye halkası",
                "Kampüs Reels",
                "Arama · önerilenler",
            ],
        ),
        (
            "Kampüs",
            [
                "Etkinlik & başvuru",
                "Topluluk hesapları",
                "Elçilik / aidiyet",
                "İş ortakları",
                "Duyuru & bildirim",
                "Çalışma odaları",
                "Deep / App Links",
            ],
        ),
        (
            "Kariyer",
            [
                "CV-AI",
                "Staj-AI",
                "Firma web paneli",
                "İlan & başvuru",
                "Mail / teklif",
                "Başvuru sıralama AI",
                "Plus · Market",
            ],
        ),
        (
            "Yönetim",
            [
                "Admin portal",
                "Şikâyet & yasak",
                "Rozet / rol",
                "Guard moderasyon",
                "KVKK / gizlilik",
                "Oturum güvenliği",
                "Hesap silme",
            ],
        ),
    ]
    for i, (title, items) in enumerate(cols):
        x = Inches(0.55 + i * 3.15)
        rounded_rect(slide, x, Inches(1.75), Inches(3.0), Inches(4.85), NAVY2 if i % 2 == 0 else SOFT, line=None if i % 2 == 0 else LINE)
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(1.95),
            Inches(2.6),
            Inches(0.35),
            title,
            size=16,
            bold=True,
            color=CYAN if i % 2 == 0 else NAVY2,
        )
        bullet_block(
            slide,
            x + Inches(0.2),
            Inches(2.45),
            Inches(2.6),
            items,
            size=12,
        )


def close(prs, blank):
    slide = prs.slides.add_slide(blank)
    dark_background(slide, prs)
    add_textbox(
        slide,
        Inches(0.85),
        Inches(1.15),
        Inches(11.5),
        Inches(0.4),
        "Saygılarımızla arz ederiz",
        size=16,
        color=CYAN,
    )
    add_textbox(
        slide,
        Inches(0.85),
        Inches(1.75),
        Inches(11.5),
        Inches(1.5),
        "GAÜN’den Türkiye’ye örnek\nbir dijital kampüs modeli",
        size=36,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(0.85),
        Inches(3.55),
        Inches(11.0),
        Inches(0.9),
        "Öğrenci deneyimi · Resmi topluluklar · Etkinlik ve SKS süreçleri · "
        "Kariyer köprüsü · Güvenli kampüs iletişimi — tek doğrulanmış platformda.",
        size=16,
        color=RGBColor(0xD5, 0xE4, 0xEE),
    )
    rounded_rect(slide, Inches(0.85), Inches(4.8), Inches(5.6), Inches(1.35), NAVY3)
    add_textbox(
        slide,
        Inches(1.1),
        Inches(5.05),
        Inches(5.1),
        Inches(0.9),
        "KampüsteyimAPP · AYS Tech\nKayra Çatalkaya\ninfo@kampusteyim.app",
        size=14,
        bold=True,
        color=WHITE,
    )
    rounded_rect(slide, Inches(6.8), Inches(4.8), Inches(5.7), Inches(1.35), TEAL)
    add_textbox(
        slide,
        Inches(7.05),
        Inches(5.1),
        Inches(5.2),
        Inches(0.85),
        "kampusteyim.app\napp.kampusteyim.app\n@kampusteyimapp",
        size=15,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )


def to_ppsx(pptx_path: Path) -> Path:
    ppsx_path = pptx_path.with_suffix(".ppsx")
    ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    ET.register_namespace("", ns)
    with zipfile.ZipFile(pptx_path, "r") as src, NamedTemporaryFile(
        suffix=".ppsx", delete=False
    ) as tmp:
        tmp_path = Path(tmp.name)
    with zipfile.ZipFile(pptx_path, "r") as src, zipfile.ZipFile(
        tmp_path, "w", zipfile.ZIP_DEFLATED
    ) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "[Content_Types].xml":
                root = ET.fromstring(data)
                for node in root:
                    if node.attrib.get("PartName") == "/ppt/presentation.xml":
                        node.set(
                            "ContentType",
                            "application/vnd.openxmlformats-officedocument."
                            "presentationml.slideshow.main+xml",
                        )
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            dst.writestr(item, data)
    shutil.move(str(tmp_path), ppsx_path)
    return ppsx_path


def main():
    builders = [
        cover,
        purpose,
        vision,
        actors,
        registration,
        feed,
        stories_reels,
        discovery,
        events,
        communities,
        career,
        notifications_plus,
        admin_moderation,
        privacy_security,
        sks,
        roadmap,
        feature_inventory,
        close,
    ]
    _PAGE["total"] = len(builders) - 2  # cover + kapanış hariç
    _PAGE["n"] = 0

    prs, blank = new_deck()
    for fn in builders:
        fn(prs, blank)

    prs.core_properties.title = (
        "KampüsteyimAPP — Gaziantep Üniversitesi Rektörlüğü Sunumu"
    )
    prs.core_properties.subject = "Resmi ürün ve pilot uygulama sunumu"
    prs.core_properties.author = "KampüsteyimAPP · AYS Tech · Kayra Çatalkaya"
    prs.save(str(OUT))
    shutil.copy2(OUT, OUT_ALT)
    shutil.copy2(OUT, ROOT / OUT.name)
    shutil.copy2(OUT, ROOT / OUT_ALT.name)
    ppsx = to_ppsx(OUT)
    shutil.copy2(ppsx, ROOT / ppsx.name)
    print(f"OK: {OUT} ({len(prs.slides)} slides)")
    print(f"OK: {OUT_ALT}")
    print(f"OK: {ppsx}")
    print(f"COPY -> {ROOT}")


if __name__ == "__main__":
    main()
