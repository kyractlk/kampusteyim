# -*- coding: utf-8 -*-
"""Firma ve üniversite toplulukları için ayrı KampüsteyimAPP sunumları."""
from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from tempfile import NamedTemporaryFile
from xml.etree import ElementTree as ET

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_pptx import (
    BRAND,
    CYAN,
    INK,
    INSTAGRAM,
    INSTAGRAM_URL,
    LINE,
    MUTED,
    NAVY,
    NAVY2,
    NAVY3,
    ROOT,
    SHOTS,
    SOFT,
    TEAL,
    WHITE,
    add_multiline,
    add_shot,
    add_textbox,
    bullet_block,
    card_with_text,
    footer_bar,
    header_bar,
    matrix_2x2,
    process_row,
    rounded_rect,
    section_title,
    stack_layers,
)

HERE = Path(__file__).resolve().parent
DOWNLOAD_URL = "https://kampusteyim.app/get.html"
INSTAGRAM_LINK = "https://instagram.com/kampusteyimapp"

FIRM_PPTX = HERE / "KampusteyimAPP_Firmalar_Sunumu.pptx"
COMMUNITY_PPTX = HERE / "KampusteyimAPP_Universite_Topluluklari_Sunumu.pptx"


def qr_image(url: str, name: str) -> Path:
    out = HERE / "shots" / name
    out.parent.mkdir(parents=True, exist_ok=True)
    qr = qrcode.QRCode(version=None, box_size=12, border=3)
    qr.add_data(url)
    qr.make(fit=True)
    qr.make_image(fill_color="#071824", back_color="white").save(out)
    return out


def add_qr(slide, image: Path, url: str, left, top, size, title: str, subtitle: str):
    rounded_rect(
        slide,
        left - Inches(0.12),
        top - Inches(0.12),
        size + Inches(0.24),
        size + Inches(1.08),
        WHITE,
        line=LINE,
    )
    pic = slide.shapes.add_picture(str(image), left, top, size, size)
    pic.click_action.hyperlink.address = url
    add_textbox(
        slide,
        left - Inches(0.05),
        top + size + Inches(0.12),
        size + Inches(0.1),
        Inches(0.3),
        title,
        size=14,
        bold=True,
        color=NAVY2,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        left - Inches(0.05),
        top + size + Inches(0.43),
        size + Inches(0.1),
        Inches(0.38),
        subtitle,
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def new_deck() -> tuple[Presentation, object]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slide_layouts[6]


def cover(prs, blank, audience: str, subtitle: str, badge: str):
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    for y, color in ((0, TEAL), (prs.slide_height - Inches(0.3), TEAL)):
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, Inches(0.3)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()
    logo = SHOTS / "logo_kampus_dark.png"
    if not logo.exists():
        logo = BRAND / "kampusteyim_icon.png"
    if logo.exists():
        slide.shapes.add_picture(
            str(logo), Inches(5.9), Inches(0.85), Inches(1.5), Inches(1.5)
        )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.55),
        Inches(11.75),
        Inches(0.7),
        "KampüsteyimAPP",
        size=40,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1),
        Inches(3.28),
        Inches(11.3),
        Inches(0.45),
        audience,
        size=23,
        bold=True,
        color=CYAN,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1.7),
        Inches(3.92),
        Inches(9.9),
        Inches(0.7),
        subtitle,
        size=15,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    rounded_rect(
        slide, Inches(4.55), Inches(5.0), Inches(4.25), Inches(0.58), NAVY3
    )
    add_textbox(
        slide,
        Inches(4.55),
        Inches(5.14),
        Inches(4.25),
        Inches(0.3),
        badge,
        size=13,
        bold=True,
        color=CYAN,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1),
        Inches(6.22),
        Inches(11.3),
        Inches(0.3),
        f"{INSTAGRAM}  ·  kampusteyim.app",
        size=12,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    return slide


def finish(prs: Presentation, out: Path, header: str) -> Path:
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        footer_bar(slide, prs, i, total - 1)
    prs.core_properties.title = header
    prs.core_properties.subject = "KampüsteyimAPP hedef kitle sunumu"
    prs.core_properties.author = "KampüsteyimAPP · AYS Tech"
    prs.save(str(out))
    shutil.copy2(out, ROOT / out.name)
    return out


def build_firm_deck() -> Path:
    prs, blank = new_deck()
    slides = []
    slides.append(
        cover(
            prs,
            blank,
            "FİRMALAR İÇİN KAMPÜS İŞ BİRLİĞİ",
            "İşveren markası, staj ve yetenek erişimini tek doğrulanmış kampüs kanalında yönetin.",
            "Firma Online · Staj-AI · Firma AI · CV-AI · Etkinlik · Reels",
        )
    )

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "NEDEN KAMPÜSTEYİMAPP?", "Öğrenci kitlesine doğrudan ve güvenli erişim")
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(1.8),
        Inches(11.9),
        Inches(4.6),
        [
            ("İşveren markası", "Doğrulanmış firma hesabıyla akış, hikâye ve Reels’te görünür olun."),
            ("Staj & ilan", "Staj-AI ile ilan, başvuru ve aday havuzunu tek panelde yönetin."),
            ("Kampüs etkinliği", "Kariyer günü ve atölyeleri kota + başvuru ile yayınlayın."),
            ("Güven & Guard", "Rozet, moderasyon ve doğrulanmış hesapla marka güvenini koruyun."),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "FİRMA ONLINE PANELİ", "Web’de güçlü işveren operasyonu")
    process_row(
        s,
        Inches(0.55),
        Inches(1.95),
        Inches(12.2),
        [
            ("Giriş", "Admin’in açtığı hesap"),
            ("İlan", "Staj / iş / part-time"),
            ("Aday", "CV tarama + sıralama"),
            ("İletişim", "Mail · direkt teklif"),
        ],
    )
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(3.55),
        Inches(11.9),
        Inches(2.85),
        [
            ("Canlı panel", "Açık / kapalı ilan, başvuru ve AI-hazır metrikleri anlık."),
            ("Zengin ilan", "Departman · ofis/hibrit/uzaktan · etiket · son başvuru."),
            ("Firma AI", "İlan seçerek başvuruları gerekçeli skorla sıralayın."),
            ("Öğrenci tarama", "Kampüs profili ve CV’si olan adaylara ulaşın."),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "FİRMA AI", "Kısa listeyi dakikalar içinde kurun")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(12.2),
        Inches(0.55),
        "Her açık ilan için başvuranları CV, yetkinlik ve gereksinim uyumuna göre skorlayın; güçlü yön / eksik gerekçesiyle kısa liste çıkarın.",
        size=14,
        color=INK,
    )
    process_row(
        s,
        Inches(0.55),
        Inches(2.4),
        Inches(12.2),
        [
            ("İlan seç", "Başvurulu açık fırsat"),
            ("Sırala", "Cloud / yerel AI skor"),
            ("İncele", "Profil · CV · gerekçe"),
            ("Teklif", "Direkt teklif + push"),
        ],
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(4.15),
        Inches(12.0),
        [
            "Birden fazla ilanda AI hazır sayacı — paneldan tek tıkla sıralama",
            "Skor kartında güçlü yönler ve eksikler görünür",
            "Teklif gönderildiğinde öğrenci Staj-AI kutusu ve bildirimi güncellenir",
            "Mail ile adaya kurumsal mesaj gönderme seçeneği",
        ],
        size=13,
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "STAJ-AI (ÖĞRENCİ)", "CV kapısından nitelikli başvuruya")
    process_row(
        s,
        Inches(0.55),
        Inches(1.95),
        Inches(12.2),
        [
            ("CV-AI", "ATS özgeçmiş"),
            ("Keşif", "Arama + filtre"),
            ("Uyum", "Beceri eşleşme skoru"),
            ("Başvuru", "Teklif kutusu"),
        ],
    )
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(3.55),
        Inches(11.9),
        Inches(2.85),
        [
            ("Kalite kapısı", "CV hazır değilse başvuru açılmaz — havuz daha nitelikli."),
            ("Akıllı liste", "Staj / iş / part-time, uzaktan-hibrit ve metin araması."),
            ("Uyum skoru", "Öğrenci becerileri ilan gereksinimleriyle örtüşürse öne çıkar."),
            ("Canlı teklif", "Firma teklifleri anlık; başvurularım paneli takip kolaylaştırır."),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "İŞVEREN MARKASI", "Akış, hikâye ve Reels ile kampüste görünür olun")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.0),
        Inches(0.65),
        "Öğrenci diline uygun içeriklerle kurum kültürünüzü, ekiplerinizi ve fırsatlarınızı anlatın.",
        size=14,
        color=INK,
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(2.4),
        Inches(7.0),
        [
            "Doğrulanmış firma rozetiyle güven verin",
            "Şehir / üniversite filtreli akışta doğru kitleye düşün",
            "Hikâye halkası ve Reels ile kısa, yüksek etkileşimli anlatım",
            "Stream-first medya: içerik anında açılır, bekleme azalır",
            "Deep link ile ilan / profil paylaşımını mağazaya kadar taşıyın",
        ],
        size=13,
    )
    add_shot(
        s, SHOTS / "03_feed.png", Inches(8.45), Inches(1.65), Inches(4.3), Inches(5.15)
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "KARİYER HUNİSİ", "İlandan nitelikli adaya")
    process_row(
        s,
        Inches(0.55),
        Inches(2.0),
        Inches(12.2),
        [
            ("Firma profili", "Doğrulanmış kimlik"),
            ("Staj-AI", "İlan ve kriterler"),
            ("Başvuru", "Öğrenci aday havuzu"),
            ("Eşleşme", "CV-AI + Firma AI"),
        ],
    )
    stack_layers(
        s,
        Inches(0.8),
        Inches(3.65),
        Inches(11.7),
        [
            ("Erişim", "Kampüs öğrencileri ve bölümler", TEAL),
            ("Nitelik", "CV-AI · profil · yetkinlik · uyum skoru", RGBColor(0x17, 0x8A, 0x84)),
            ("Operasyon", "Başvuru · AI sıralama · mail · teklif", NAVY2),
            ("Güven", "Doğrulanmış hesap · Guard · push", NAVY),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "ETKİNLİK OPERASYONU", "Kariyer günü ve atölyeleri ölçülebilir yönetin")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.1),
        Inches(0.65),
        "Tarih, yer, hedef kitle, son başvuru ve kontenjan tek kartta; başvurular canlı görünür.",
        size=14,
        color=INK,
    )
    process_row(
        s,
        Inches(0.55),
        Inches(2.5),
        Inches(7.2),
        [
            ("Oluştur", "Etkinlik kartı"),
            ("Duyur", "Akış + hikâye + Reels"),
            ("Başvuru", "Öğrenci katılımı"),
            ("Ölç", "Kota ve ilgi"),
        ],
    )
    add_shot(
        s,
        SHOTS / "06_events.png",
        Inches(8.55),
        Inches(1.65),
        Inches(4.15),
        Inches(5.1),
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "GÖRÜNÜRLÜK & PLUS", "Kampüste öne çıkmanın yeni katmanları")
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(1.85),
        Inches(11.9),
        Inches(4.5),
        [
            ("KampüsteyimPlus", "Öğrenci kitlesinde yeşil tick ve öncelikli deneyim."),
            ("Market", "app.kampusteyim.app/market — web’den Plus satışı."),
            ("İş ortaklarımız", "Kurumsal partner vitrini ile güven ve görünürlük."),
            ("Anlık bildirim", "İlan, teklif ve etkileşim push + uygulama içi inbox."),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "AKTİVASYON", "Firmanız 4 adımda kampüste")
    process_row(
        s,
        Inches(0.55),
        Inches(2.1),
        Inches(12.2),
        [
            ("Başvuru", "Firma hesabı talebi"),
            ("Doğrulama", "Yetkili ve kurum kontrolü"),
            ("İlk yayın", "Marka / ilan / etkinlik"),
            ("Büyüme", "Etkileşim ve başvuru"),
        ],
    )
    rounded_rect(s, Inches(1.1), Inches(4.2), Inches(11.1), Inches(1.2), SOFT, line=TEAL)
    add_textbox(
        s,
        Inches(1.45),
        Inches(4.48),
        Inches(10.4),
        Inches(0.55),
        "Firma hesapları yalnızca platform admini tarafından açılır ve teslim edilir.",
        size=17,
        bold=True,
        color=NAVY2,
        align=PP_ALIGN.CENTER,
    )
    slides.append(s)

    download_qr = qr_image(DOWNLOAD_URL, "qr_kampusteyim_indir.png")
    insta_qr = qr_image(INSTAGRAM_LINK, "qr_kampusteyim_instagram.png")
    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Firmalar")
    section_title(s, "İLETİŞİM", "Kampüse birlikte değer katalım")
    add_qr(
        s,
        download_qr,
        DOWNLOAD_URL,
        Inches(2.25),
        Inches(1.85),
        Inches(2.45),
        "Uygulamayı incele",
        "iOS / Android akıllı yönlendirme",
    )
    add_qr(
        s,
        insta_qr,
        INSTAGRAM_LINK,
        Inches(8.6),
        Inches(1.85),
        Inches(2.45),
        "Hesabı takip et",
        INSTAGRAM,
    )
    add_textbox(
        s,
        Inches(1.2),
        Inches(5.55),
        Inches(10.9),
        Inches(0.7),
        "KampüsteyimAPP · AYS Tech · info@kampusteyim.app\nkampusteyim.app  ·  app.kampusteyim.app  ·  /market · /firma",
        size=15,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    slides.append(s)

    return finish(prs, FIRM_PPTX, "KampüsteyimAPP — Firmalar Sunumu")


def build_community_deck() -> Path:
    prs, blank = new_deck()
    slides = []
    slides.append(
        cover(
            prs,
            blank,
            "ÜNİVERSİTE TOPLULUKLARI İÇİN",
            "Resmi kimliğinizi güçlendirin; duyuru, etkinlik ve üye etkileşimini kampüsün dijital alanında büyütün.",
            "Rozet · Hikâye · Reels · Etkinlik · Elçilik",
        )
    )

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "TOPLULUK DEĞERİ", "Resmi topluluğunuz için tek dijital merkez")
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(1.8),
        Inches(11.9),
        Inches(4.6),
        [
            ("Resmi kimlik", "Altın / mavi rozetle sahte hesaplardan net biçimde ayrışın."),
            ("Üye kazanımı", "Standdan QR ile uygulamaya ve topluluk hesabınıza yönlendirin."),
            ("Etkinlik yönetimi", "Tarih, yer, kota ve başvuruyu tek karttan yönetin."),
            ("Canlı etkileşim", "Akış, hikâye halkası ve Reels ile nabzı canlı tutun."),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "STAND AKIŞI", "Tanıtım standında 30 saniyelik dönüşüm")
    process_row(
        s,
        Inches(0.55),
        Inches(1.9),
        Inches(12.2),
        [
            ("QR okut", "Uygulamayı aç / indir"),
            ("Hesabı takip et", INSTAGRAM),
            ("Topluluğu bul", "Resmi hesap / rozet"),
            ("Etkinliğe katıl", "Başvuru ve duyuru"),
        ],
    )
    rounded_rect(s, Inches(0.8), Inches(3.65), Inches(11.7), Inches(1.55), NAVY2)
    add_textbox(
        s,
        Inches(1.1),
        Inches(3.98),
        Inches(11.1),
        Inches(0.45),
        "Stand çağrısı",
        size=17,
        bold=True,
        color=CYAN,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(1.1),
        Inches(4.45),
        Inches(11.1),
        Inches(0.45),
        "“KampüsteyimAPP hesabımızı takip et, QR’ı okut ve topluluk duyurularını kaçırma.”",
        size=18,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "İÇERİK PLANI", "Topluluğunuz kampüste sürekli görünür")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.0),
        Inches(0.55),
        "Tek bir duyuru yerine etkinlik öncesi, anı ve sonrasını kapsayan bir iletişim ritmi kurun.",
        size=14,
        color=INK,
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(2.35),
        Inches(7.0),
        [
            "Akış: ana duyuru — şehir / üniversite filtresiyle doğru kitle",
            "Hikâye halkası: geri sayım, hatırlatma, anlık paylaşım (stream)",
            "Reels: stand, ekip ve etkinlik özeti — kaydırınca anında oynar",
            "Profil: logo, rozet, biyografi ve geçmiş içerikler",
            "Push: takipçiye anlık duyuru ve etkinlik hatırlatması",
        ],
        size=13,
    )
    add_shot(
        s, SHOTS / "03_feed.png", Inches(8.5), Inches(1.65), Inches(4.25), Inches(5.1)
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "HİKÂYE & REELS", "Instagram ritminde kampüs medyası")
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(1.85),
        Inches(11.9),
        Inches(4.5),
        [
            ("Hikâye halkası", "Görülmemiş / görülmüş halka; canlı snapshot ile saniyelik yansıma."),
            ("Stream-first", "Video tam indirme beklemez — progressive oynatma."),
            ("Reels motoru", "Önceki / sonraki ısıtma; kaydırınca takılma azalır."),
            ("Deep link", "Paylaşılan hikâye / profil / etkinlik mağazaya kadar gider."),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "ETKİNLİK", "Başvuru, kontenjan ve hatırlatma tek yerde")
    process_row(
        s,
        Inches(0.55),
        Inches(2.0),
        Inches(7.2),
        [
            ("Oluştur", "Tarih · yer · görsel"),
            ("Yayınla", "Akış + hikâye + Reels"),
            ("Topla", "Başvuru + kota"),
            ("Yönet", "Katılım takibi"),
        ],
    )
    add_textbox(
        s,
        Inches(0.55),
        Inches(3.55),
        Inches(7.1),
        Inches(0.4),
        "Topluluk için kazanım",
        size=16,
        bold=True,
        color=NAVY2,
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(4.0),
        Inches(7.0),
        [
            "Dağınık form ve mesaj trafiği azalır",
            "Kontenjan ve son başvuru net görünür",
            "Üyeler duyuruyu kampüs akışında keşfeder",
            "App Link ile etkinlik linki uygulama yoksa mağazaya düşer",
        ],
        size=13,
    )
    add_shot(
        s,
        SHOTS / "06_events.png",
        Inches(8.55),
        Inches(1.65),
        Inches(4.15),
        Inches(5.1),
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "ELÇİLİK & AİDİYET", "Kampüs elçileri ve bağlı hesaplar")
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(1.85),
        Inches(11.9),
        Inches(4.5),
        [
            ("Kampüs elçisi", "Topluluğu temsil eden öğrencilere özel rozet ve görünürlük."),
            ("Bağlı hesaplar", "Öğrenci profilinde topluluk / firma aidiyeti gösterimi."),
            ("Takip & gizlilik", "Gizli hesap ve takip isteği ile üye yönetimini netleştirin."),
            ("İş ortakları", "About bölümünde partner vitrini — sponsor görünürlüğü."),
        ],
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "RESMİ HESAP", "Topluluk hesabınızı aktive edin")
    process_row(
        s,
        Inches(0.55),
        Inches(2.0),
        Inches(12.2),
        [
            ("Başvuru", "Topluluk adı ve yetkili"),
            ("Doğrulama", "Üniversite / danışman"),
            ("Rozet", "Resmi hesap kimliği"),
            ("Yayın", "İlk etkinlik ve duyuru"),
        ],
    )
    rounded_rect(s, Inches(0.9), Inches(3.8), Inches(11.55), Inches(1.5), SOFT, line=TEAL)
    add_textbox(
        s,
        Inches(1.2),
        Inches(4.1),
        Inches(10.95),
        Inches(0.38),
        "Doğrulanmış topluluk hesabı",
        size=18,
        bold=True,
        color=NAVY2,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(1.2),
        Inches(4.55),
        Inches(10.95),
        Inches(0.5),
        "Logo zorunlu · Duyuru yetkisi · Etkinlik kartı · Guard / moderasyon",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    slides.append(s)

    download_qr = qr_image(DOWNLOAD_URL, "qr_kampusteyim_indir.png")
    insta_qr = qr_image(INSTAGRAM_LINK, "qr_kampusteyim_instagram.png")
    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "STAND QR'LARI", "Okut, indir ve takip et")
    add_qr(
        s,
        download_qr,
        DOWNLOAD_URL,
        Inches(2.15),
        Inches(1.75),
        Inches(2.65),
        "KampüsteyimAPP’i aç",
        "iOS / Android akıllı yönlendirme",
    )
    add_qr(
        s,
        insta_qr,
        INSTAGRAM_LINK,
        Inches(8.5),
        Inches(1.75),
        Inches(2.65),
        "Hesabımızı takip et",
        INSTAGRAM,
    )
    add_textbox(
        s,
        Inches(1.1),
        Inches(5.75),
        Inches(11.1),
        Inches(0.48),
        "Standda göster · Çıktı al · Öğrenciyi tek adımda yönlendir",
        size=17,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    slides.append(s)

    s = prs.slides.add_slide(blank)
    header_bar(s, prs, "Üniversite Toplulukları")
    section_title(s, "KAPANIŞ", "Topluluğunuz kampüsün dijital akışında")
    rounded_rect(s, Inches(0.9), Inches(1.9), Inches(11.55), Inches(2.35), NAVY2)
    add_textbox(
        s,
        Inches(1.3),
        Inches(2.3),
        Inches(10.75),
        Inches(0.55),
        "Etkinliğini duyur. Üyeni büyüt. Resmi kimliğini güçlendir.",
        size=24,
        bold=True,
        color=CYAN,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(1.7),
        Inches(3.08),
        Inches(9.95),
        Inches(0.65),
        "KampüsteyimAPP ile topluluğunuzu doğrulanmış, erişilebilir ve ölçülebilir bir dijital merkeze taşıyın.",
        size=16,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(1),
        Inches(5.0),
        Inches(11.3),
        Inches(0.55),
        f"{INSTAGRAM}  ·  kampusteyim.app  ·  app.kampusteyim.app",
        size=17,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    slides.append(s)

    return finish(
        prs,
        COMMUNITY_PPTX,
        "KampüsteyimAPP — Üniversite Toplulukları Sunumu",
    )


def to_ppsx(pptx_path: Path) -> Path:
    """PPTX paketini PowerPoint Show (.ppsx) içerik türüne dönüştür."""
    ppsx_path = pptx_path.with_suffix(".ppsx")
    ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    ET.register_namespace("", ns)
    with zipfile.ZipFile(pptx_path, "r") as src, NamedTemporaryFile(
        suffix=".ppsx", delete=False
    ) as tmp:
        tmp_path = Path(tmp.name)
    with zipfile.ZipFile(pptx_path, "r") as src, zipfile.ZipFile(
        tmp_path, "w", zipfile.ZIP_DEFLATED
    ) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "[Content_Types].xml":
                root = ET.fromstring(data)
                for node in root:
                    if node.attrib.get("PartName") == "/ppt/presentation.xml":
                        node.set(
                            "ContentType",
                            "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml",
                        )
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            dst.writestr(item, data)
    shutil.move(str(tmp_path), ppsx_path)
    shutil.copy2(ppsx_path, ROOT / ppsx_path.name)
    return ppsx_path


def main():
    firm = build_firm_deck()
    community = build_community_deck()
    outputs = [firm, to_ppsx(firm), community, to_ppsx(community)]
    corporate = HERE / "KampusteyimAPP_Kurumsal_Tanitim.pptx"
    if corporate.exists():
        outputs.append(to_ppsx(corporate))
    for output in outputs:
        print(output)


if __name__ == "__main__":
    main()
