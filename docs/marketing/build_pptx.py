# -*- coding: utf-8 -*-
"""KampüsteyimAPP kurumsal tanıtım — PowerPoint (16:9)."""
from __future__ import annotations

from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = Path(__file__).resolve().parent / "KampusteyimAPP_Kurumsal_Tanitim.pptx"
SHOTS = Path(__file__).resolve().parent / "shots"
BRAND = ROOT / "assets" / "brand"
LOGOS = ROOT / "assets" / "logos"

INSTAGRAM = "@kampusteyimapp"
INSTAGRAM_URL = "instagram.com/kampusteyimapp"

NAVY = RGBColor(0x07, 0x18, 0x24)
NAVY2 = RGBColor(0x0B, 0x1F, 0x33)
NAVY3 = RGBColor(0x12, 0x32, 0x4A)
TEAL = RGBColor(0x1F, 0xA6, 0xA0)
CYAN = RGBColor(0x2E, 0xC4, 0xB6)
SOFT = RGBColor(0xEE, 0xF5, 0xF8)
INK = RGBColor(0x1A, 0x2B, 0x3A)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xD8, 0xE2)


def _set_run(run, size=14, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_multiline(slide, left, top, width, height, lines, *, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT, gap_pt=4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap_pt)
        run = p.add_run()
        run.text = line
        _set_run(run, size=size, bold=bold, color=color)
    return box


def rounded_rect(slide, left, top, width, height, fill, *, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    # milder corner
    try:
        shape.adjustments[0] = 0.12
    except Exception:
        pass
    return shape


def pill(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    return shape


def card_with_text(slide, left, top, width, height, title, subtitle_lines, fill=NAVY2, title_color=WHITE, sub_color=None):
    if sub_color is None:
        sub_color = RGBColor(0xC9, 0xDB, 0xE8)
    rounded_rect(slide, left, top, width, height, fill)
    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.18),
        width - Inches(0.36),
        Inches(0.35),
        title,
        size=16,
        bold=True,
        color=title_color,
        align=PP_ALIGN.CENTER,
    )
    add_multiline(
        slide,
        left + Inches(0.14),
        top + Inches(0.52),
        width - Inches(0.28),
        height - Inches(0.6),
        subtitle_lines,
        size=11,
        color=sub_color,
        align=PP_ALIGN.CENTER,
        gap_pt=2,
    )


def connector(slide, x1, y1, x2, y2, color=TEAL):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    return line


def header_bar(slide, prs, title_right="Kurumsal Tanıtım"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY2
    bar.line.fill.background()
    logo = SHOTS / "logo_kampus_header.png"
    if not logo.exists():
        logo = BRAND / "kampusteyim_icon.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.35), Inches(0.1), Inches(0.35), Inches(0.35))
    add_textbox(slide, Inches(0.8), Inches(0.12), Inches(4), Inches(0.35), "KampüsteyimAPP", size=14, bold=True, color=WHITE)
    add_textbox(
        slide,
        prs.slide_width - Inches(3.2),
        Inches(0.14),
        Inches(2.8),
        Inches(0.3),
        title_right,
        size=11,
        color=RGBColor(0xB8, 0xC8, 0xD4),
        align=PP_ALIGN.RIGHT,
    )


def footer_bar(slide, prs, page: int, total: int):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.42), prs.slide_width, Inches(0.42)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    add_textbox(slide, Inches(0.4), prs.slide_height - Inches(0.34), Inches(4), Inches(0.28), INSTAGRAM, size=11, bold=True, color=WHITE)
    add_textbox(
        slide,
        prs.slide_width - Inches(2.2),
        prs.slide_height - Inches(0.34),
        Inches(1.8),
        Inches(0.28),
        f"{page} / {total}",
        size=11,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
    )


def section_title(slide, kicker: str, title: str):
    add_textbox(slide, Inches(0.55), Inches(0.72), Inches(4), Inches(0.28), kicker, size=11, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.55), Inches(0.98), Inches(11.5), Inches(0.45), title, size=28, bold=True, color=NAVY2)
    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.45), Inches(1.1), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def process_row(slide, left, top, width, steps: list[tuple[str, str]]):
    n = len(steps)
    gap = Inches(0.18)
    box_w = (width - gap * (n - 1)) / n
    box_h = Inches(1.15)
    for i, (title, sub) in enumerate(steps):
        x = left + i * (box_w + gap)
        rounded_rect(slide, x, top, box_w, box_h, SOFT if i % 2 == 0 else RGBColor(0xE4, 0xF0, 0xF4))
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.14), top + Inches(0.18), Inches(0.32), Inches(0.32))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        add_textbox(slide, x + Inches(0.14), top + Inches(0.2), Inches(0.32), Inches(0.3), str(i + 1), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.55), top + Inches(0.18), box_w - Inches(0.7), Inches(0.32), title, size=13, bold=True, color=NAVY2)
        add_textbox(slide, x + Inches(0.55), top + Inches(0.55), box_w - Inches(0.7), Inches(0.45), sub, size=11, color=MUTED)
        if i < n - 1:
            connector(
                slide,
                x + box_w,
                top + box_h / 2,
                x + box_w + gap,
                top + box_h / 2,
            )


def stack_layers(slide, left, top, width, layers: list[tuple[str, str, RGBColor]]):
    h = Inches(0.58)
    gap = Inches(0.1)
    for i, (title, sub, fill) in enumerate(layers):
        y = top + i * (h + gap)
        rounded_rect(slide, left, y, width, h, fill)
        light = fill in (NAVY, NAVY2, TEAL, RGBColor(0x17, 0x8A, 0x84))
        tc = WHITE if light else NAVY2
        sc = RGBColor(0xD5, 0xE4, 0xEE) if light else MUTED
        add_textbox(slide, left + Inches(0.25), y + Inches(0.14), Inches(5.5), Inches(0.35), title, size=14, bold=True, color=tc)
        add_textbox(slide, left + width - Inches(5.2), y + Inches(0.14), Inches(4.9), Inches(0.35), sub, size=12, color=sc, align=PP_ALIGN.RIGHT)


def matrix_2x2(slide, left, top, width, height, cells: list[tuple[str, str]]):
    gap = Inches(0.16)
    cw = (width - gap) / 2
    ch = (height - gap) / 2
    fills = [NAVY2, NAVY3, RGBColor(0x0F, 0x3A, 0x3A), RGBColor(0x1A, 0x33, 0x48)]
    for i, (title, body) in enumerate(cells[:4]):
        col, row = i % 2, i // 2
        x = left + col * (cw + gap)
        y = top + row * (ch + gap)
        rounded_rect(slide, x, y, cw, ch, fills[i])
        add_textbox(slide, x + Inches(0.22), y + Inches(0.2), cw - Inches(0.4), Inches(0.35), title, size=15, bold=True, color=CYAN)
        add_textbox(slide, x + Inches(0.22), y + Inches(0.55), cw - Inches(0.4), ch - Inches(0.7), body, size=12, color=RGBColor(0xC9, 0xDB, 0xE8))


def bullet_block(slide, left, top, width, items: list[str], *, size=13):
    box = slide.shapes.add_textbox(left, top, width, Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        _set_run(run, size=size, color=INK)


def add_shot(slide, path: Path, left, top, max_w, max_h, *, center=True):
    """Telefon ekran görüntüsünü en-boy oranını bozmadan yerleştir."""
    if not path.exists():
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    if iw <= 0 or ih <= 0:
        return None
    scale = min(float(max_w) / iw, float(max_h) / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    x = left
    if center:
        x = left + (max_w - w) // 2
    return slide.shapes.add_picture(str(path), x, top, w, h)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides_meta: list = []

    # ── Kapak ──
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    top_strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.28))
    top_strip.fill.solid()
    top_strip.fill.fore_color.rgb = TEAL
    top_strip.line.fill.background()
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.35), prs.slide_width, Inches(0.35))
    bot.fill.solid()
    bot.fill.fore_color.rgb = TEAL
    bot.line.fill.background()

    logo = SHOTS / "logo_kampus_dark.png"
    if not logo.exists():
        logo = BRAND / "kampusteyim_icon.png"
    if logo.exists():
        s.shapes.add_picture(str(logo), Inches(5.9), Inches(1.1), Inches(1.5), Inches(1.5))
    add_textbox(s, Inches(1), Inches(2.8), Inches(11.3), Inches(0.7), "KampüsteyimAPP", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.4), "Kampüsün kurumsal dijital ekosistemi", size=18, color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(4.15), Inches(11.3), Inches(0.35), "FİRMALAR  ·  RESMİ ÜNİVERSİTE TOPLULUKLARI", size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        Inches(2),
        Inches(4.65),
        Inches(9.3),
        Inches(0.7),
        "Ekosistem şeması · operasyon süreçleri · kariyer hunisi\nGüvenli, doğrulanmış ve kampüse özel kurumsal sunum",
        size=14,
        color=RGBColor(0xB8, 0xE8, 0xE3),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(s, Inches(1), Inches(5.7), Inches(11.3), Inches(0.3), f"Instagram  {INSTAGRAM}  ·  {INSTAGRAM_URL}", size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(6.15), Inches(11.3), Inches(0.3), "Kayra Çatalkaya  ·  AYS Tech  ·  Temmuz 2026", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    slides_meta.append(s)

    # ── 01 Ekosistem ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 01", "Kampüs ekosistemi")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(12.2),
        Inches(0.7),
        "KampüsteyimAPP; öğrenci, firma ve resmi üniversite topluluklarını tek doğrulanmış dijital çatı altında birleştirir. "
        "Amaç dağınık kanallar değil — ölçülebilir, güvenli ve kampüse özel bir etkileşim ağıdır.",
        size=14,
        color=INK,
    )

    # Üç aktör — eşit aralıklı
    actors = [
        ("Öğrenci", ["Sosyal · CV-AI", "Başvuru · Odak"]),
        ("Firma", ["İşveren markası", "Staj-AI · İlan"]),
        ("Topluluk", ["Resmi kimlik", "Etkinlik · Duyuru"]),
    ]
    aw, ah = Inches(3.2), Inches(1.35)
    gap = Inches(0.45)
    total_w = 3 * aw + 2 * gap
    start_x = (prs.slide_width - total_w) / 2
    actor_y = Inches(2.55)
    actor_centers = []
    for i, (title, subs) in enumerate(actors):
        x = start_x + i * (aw + gap)
        card_with_text(s, x, actor_y, aw, ah, title, subs, fill=NAVY3)
        actor_centers.append(x + aw / 2)

    # Merkez platform + temiz T-şeklinde bağlantı
    cw, ch = Inches(5.2), Inches(1.45)
    cx = (prs.slide_width - cw) / 2
    cy = Inches(4.7)
    mid_y = actor_y + ah + Inches(0.35)
    for acx in actor_centers:
        connector(s, acx, actor_y + ah, acx, mid_y)
    connector(s, actor_centers[0], mid_y, actor_centers[-1], mid_y)
    platform_cx = cx + cw / 2
    connector(s, platform_cx, mid_y, platform_cx, cy)
    card_with_text(
        s,
        cx,
        cy,
        cw,
        ch,
        "KampüsteyimAPP",
        ["Doğrulama · Akış · Kariyer", "Etkinlik · Guard · Plus"],
        fill=NAVY2,
    )
    add_textbox(
        s,
        Inches(0.55),
        Inches(6.35),
        Inches(12.2),
        Inches(0.3),
        "Tek kampüs ekosistemi · doğrulanmış hesaplar · ölçülebilir etkileşim",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    slides_meta.append(s)

    # ── 01b Değer matrisi ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 01", "Kurumsal değer alanları")
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(1.8),
        Inches(11.9),
        Inches(4.6),
        [
            ("Güven", "Altın / mavi / yeşil rozetler ve AYS Tech Guard ile itibar katmanı."),
            ("Erişim", "Doğrudan öğrenci kitlesi — dağınık sosyal kanallar yerine tek kampüs."),
            ("Operasyon", "Etkinlik, duyuru, kota ve başvuru tek panelde yönetilir."),
            ("Kariyer", "CV-AI + Staj-AI ile işveren–aday köprüsü hızlanır."),
        ],
    )
    slides_meta.append(s)

    # ── 02 Süreç + stack ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 02", "Kurumsal katılım süreci")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(12.2),
        Inches(0.45),
        "Firma veya resmi topluluk hesabı, onay sonrası kampüste görünür hale gelir.",
        size=14,
        color=INK,
    )
    process_row(
        s,
        Inches(0.55),
        Inches(2.25),
        Inches(12.2),
        [
            ("Başvuru", "Resmi hesap talebi"),
            ("Doğrulama", "Kimlik & rozet"),
            ("Yayın", "Akış / etkinlik"),
            ("Ölçüm", "Etkileşim & başvuru"),
        ],
    )
    add_textbox(s, Inches(0.55), Inches(3.7), Inches(12.2), Inches(0.35), "Güven ve altyapı katmanları", size=16, bold=True, color=NAVY2)
    stack_layers(
        s,
        Inches(0.55),
        Inches(4.15),
        Inches(12.2),
        [
            ("Deneyim katmanı", "Akış · Hikâye · Reels · Etkinlik", TEAL),
            ("Kariyer katmanı", "CV-AI · Staj-AI · Plus", RGBColor(0x17, 0x8A, 0x84)),
            ("Kimlik katmanı", "Doğrulama rozetleri · Gizlilik", NAVY2),
            ("Güvenlik + altyapı", "AYS Tech Guard · Firebase · iOS/Android", NAVY),
        ],
    )
    slides_meta.append(s)

    # ── 03 Giriş ekranları ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 03", "Marka deneyimi ve güvenli giriş")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(12.2),
        Inches(0.5),
        "Sade, profesyonel ve kampüs kimliğini yansıtan bir karşılama. Misafir akışa göz atabilir; kurumsal ve öğrenci hesapları güvenle giriş yapar.",
        size=14,
        color=INK,
    )
    add_shot(s, SHOTS / "01_login.png", Inches(1.6), Inches(2.15), Inches(4.4), Inches(4.5))
    add_shot(s, SHOTS / "02_login_form.png", Inches(7.3), Inches(2.15), Inches(4.4), Inches(4.5))
    add_textbox(s, Inches(1.6), Inches(6.7), Inches(4.4), Inches(0.28), "Karşılama & marka", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.3), Inches(6.7), Inches(4.4), Inches(0.28), "Kampüs hesabıyla giriş", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    slides_meta.append(s)

    # ── 04 Akış ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 04", "Kampüs akışı")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.2),
        Inches(0.7),
        "Akış; firmaların, toplulukların ve öğrencilerin aynı dilde konuştuğu canlı vitrindir. Resmi hesaplar altın / mavi rozetle, Plus üyeleri yeşil tick ile öne çıkar.",
        size=14,
        color=INK,
    )
    add_textbox(s, Inches(0.55), Inches(2.45), Inches(7.2), Inches(0.3), "Firmalar için", size=15, bold=True, color=NAVY2)
    bullet_block(
        s,
        Inches(0.55),
        Inches(2.8),
        Inches(7.0),
        [
            "Marka ve işveren hikâyesini öğrenci dilinde paylaşın",
            "Staj / kariyer içeriklerini hashtag ile keşfettirin",
            "Doğrulanmış firma hesabıyla güven inşa edin",
        ],
    )
    add_textbox(s, Inches(0.55), Inches(4.4), Inches(7.2), Inches(0.3), "Topluluklar için", size=15, bold=True, color=NAVY2)
    bullet_block(
        s,
        Inches(0.55),
        Inches(4.75),
        Inches(7.0),
        [
            "Etkinlik heyecanını ve duyuruları akışta canlı tutun",
            "Altın tick ile resmi topluluk kimliğini gösterin",
        ],
    )
    add_shot(s, SHOTS / "03_feed.png", Inches(8.55), Inches(1.75), Inches(4.2), Inches(5.0))
    slides_meta.append(s)

    # ── 05 Hikâye süreci ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 05", "Hikâye: 24 saatlik kampüs nabzı")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.0),
        Inches(0.6),
        "Anlık görünürlük için hikâyeler 24 saat sonra otomatik silinir. Etkinlik countdown’u, staj hatırlatması veya topluluk anları doğal ritimle dolaşır.",
        size=14,
        color=INK,
    )
    process_row(
        s,
        Inches(0.55),
        Inches(2.45),
        Inches(7.2),
        [
            ("Çekim", "Kamera / galeri"),
            ("Yayın", "24 saat hikâye"),
            ("İzleme", "Görüntüleyenler"),
            ("Etkileşim", "Beğeni kalbi"),
        ],
    )
    add_shot(s, SHOTS / "04_story.png", Inches(8.55), Inches(1.75), Inches(4.2), Inches(5.0))
    slides_meta.append(s)

    # ── 06 Medya stack ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 06", "Reels, medya ve içerik koruması")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(12.2),
        Inches(0.55),
        "Kampüs Reels ile kısa, dikey ve yüksek etkileşimli videolar. Fotoğraf ve videoda indirme yoktur. Dosyalar Plus ile paylaşılır.",
        size=14,
        color=INK,
    )
    stack_layers(
        s,
        Inches(0.9),
        Inches(2.5),
        Inches(11.5),
        [
            ("Reels & sosyal medya", "Beğeni · yorum · mention · bildirim", TEAL),
            ("Medya koruması", "Foto / video: izle, indirme yok", NAVY2),
            ("Dosya (Plus)", "PDF / not paylaşımı + indirme", RGBColor(0x17, 0x8A, 0x84)),
            ("AYS Tech Guard", "Metin · link · medya · dosya denetimi", NAVY),
        ],
    )
    slides_meta.append(s)

    # ── 07 Keşif ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 07", "Keşif: arama ve doğrulanmış profiller")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.0),
        Inches(0.7),
        "Kişi, @handle ve #gönderi ile kampüsü keşfedin. Firma, Topluluk, AI Bot ve üniversite etiketleri net ayrılır.",
        size=14,
        color=INK,
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(2.5),
        Inches(7.0),
        [
            "Trend hashtag’ler etkileşimle listelenir",
            "Altın tick’li resmi hesaplar ilk bakışta güven verir",
            "İşveren ve topluluk keşfi tek arama deneyiminde",
        ],
    )
    add_shot(s, SHOTS / "05_search.png", Inches(8.55), Inches(1.75), Inches(4.2), Inches(5.0))
    slides_meta.append(s)

    # ── 08 Etkinlik ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 08", "Etkinlik operasyon şeması")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.0),
        Inches(0.55),
        "Görsel, tarih, yer, hedef kitle, son başvuru ve kota tek kartta. Öğrenci Başvur der; organizatör doluluğu canlı izler.",
        size=14,
        color=INK,
    )
    process_row(
        s,
        Inches(0.55),
        Inches(2.4),
        Inches(7.2),
        [
            ("Oluştur", "Kart + kota"),
            ("Duyur", "Akış / hikâye"),
            ("Başvuru", "Öğrenci tıklar"),
            ("Yönet", "Doluluk izle"),
        ],
    )
    add_shot(s, SHOTS / "06_events.png", Inches(8.55), Inches(1.75), Inches(4.2), Inches(5.0))
    slides_meta.append(s)

    # ── 09 Kariyer hunisi ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 09", "Kariyer hunisi: CV-AI → Staj-AI")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(7.0),
        Inches(0.55),
        "Profil; vitrin, ağ ve kariyer araçlarının buluşma noktasıdır. CV-AI ATS uyumlu özgeçmiş üretir; Staj-AI ilan ve başvuruyu tek panelde toplar.",
        size=14,
        color=INK,
    )
    process_row(
        s,
        Inches(0.55),
        Inches(2.4),
        Inches(7.2),
        [
            ("Profil", "Kimlik + rozet"),
            ("CV-AI", "ATS · dil · tema"),
            ("Staj-AI", "İlan / teklif"),
            ("Eşleşme", "Başvuru hunisi"),
        ],
    )
    bullet_block(
        s,
        Inches(0.55),
        Inches(3.9),
        Inches(7.0),
        [
            "Çalışma odası: senkron sayaç + sohbet",
            "Plus: yeşil tick, dosya, CV güçleri",
            "Gizlilik: görünürlük, arama, engelleme",
        ],
    )
    add_shot(s, SHOTS / "07_profile.png", Inches(8.55), Inches(1.75), Inches(4.2), Inches(5.0))
    slides_meta.append(s)

    # ── 10 Firma / Topluluk ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 10", "Firma & topluluk değer şeması")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(12.2),
        Inches(0.4),
        "İki kurumsal aktör aynı platformda farklı operasyonel kazanımlar elde eder.",
        size=14,
        color=INK,
    )
    # Sol / sağ kolonlar — üst üste binmeden
    left_x, right_x = Inches(0.7), Inches(6.95)
    col_w, col_h = Inches(5.7), Inches(4.4)
    rounded_rect(s, left_x, Inches(2.2), col_w, col_h, NAVY2)
    rounded_rect(s, right_x, Inches(2.2), col_w, col_h, RGBColor(0x0F, 0x2F, 0x2E))
    add_textbox(s, left_x + Inches(0.3), Inches(2.4), col_w - Inches(0.6), Inches(0.4), "Firma", size=18, bold=True, color=CYAN)
    add_textbox(s, right_x + Inches(0.3), Inches(2.4), col_w - Inches(0.6), Inches(0.4), "Resmi Topluluk", size=18, bold=True, color=CYAN)
    firma = [
        "İşveren görünürlüğü",
        "Staj & işe alım hunisi",
        "Etkinlik / kariyer günü",
        "Doğrulanmış itibar",
        "Plus: dosya + yeşil tick",
        "Guard ile güvenli marka alanı",
    ]
    topluluk = [
        "Altın tick ile resmi kimlik",
        "Etkinlik: başvuru + kota",
        "Duyuru otoritesi",
        "Üye etkileşimi (hikâye/Reels)",
        "Çalışma odası / ortak odak",
        "Sahte hesaptan net ayrışma",
    ]
    for i, t in enumerate(firma):
        add_textbox(s, left_x + Inches(0.4), Inches(3.0) + Inches(0.42) * i, col_w - Inches(0.7), Inches(0.38), f"●  {t}", size=14, color=WHITE)
    for i, t in enumerate(topluluk):
        add_textbox(s, right_x + Inches(0.4), Inches(3.0) + Inches(0.42) * i, col_w - Inches(0.7), Inches(0.38), f"●  {t}", size=14, color=WHITE)
    slides_meta.append(s)

    # ── 11 Yetenek haritası ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 11", "Yetenek haritası")
    matrix_2x2(
        s,
        Inches(0.7),
        Inches(1.75),
        Inches(11.9),
        Inches(4.7),
        [
            ("Sosyal", "Akış, hikâye, Reels, mention, hashtag, beğeni, yorum, repost."),
            ("Kurumsal", "Etkinlik, duyuru, rozetler, resmi hesap, keşfet."),
            ("Kariyer", "CV-AI, Staj-AI, Plus kota/dil/tema, dosya paylaşımı."),
            ("Güven", "Guard, gizlilik, engelleme, şikâyet, medya koruması."),
        ],
    )
    slides_meta.append(s)

    # ── 12 Kapanış ──
    s = prs.slides.add_slide(blank)
    header_bar(s, prs)
    section_title(s, "BÖLÜM 12", "Birlikte büyüyelim")
    add_textbox(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(12.2),
        Inches(0.7),
        "KampüsteyimAPP; firmanızın yetenek havuzuna, topluluğunuzun üyelerine ve kampüsün nabzına aynı anda dokunmanızı sağlar.",
        size=15,
        color=INK,
    )
    process_row(
        s,
        Inches(0.55),
        Inches(2.5),
        Inches(12.2),
        [
            ("Tanışma", "Demo / sunum"),
            ("Aktivasyon", "Resmi hesap"),
            ("Yayın", "İlk içerik"),
            ("Büyüme", "Ölçüm & Plus"),
        ],
    )
    rounded_rect(s, Inches(0.55), Inches(4.0), Inches(12.2), Inches(0.85), SOFT)
    add_textbox(
        s,
        Inches(0.8),
        Inches(4.2),
        Inches(11.7),
        Inches(0.5),
        f"Instagram  {INSTAGRAM}  ·  {INSTAGRAM_URL}",
        size=16,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    # imza kutuları
    rounded_rect(s, Inches(1.5), Inches(5.15), Inches(4.5), Inches(1.15), SOFT, line=TEAL)
    rounded_rect(s, Inches(7.3), Inches(5.15), Inches(4.5), Inches(1.15), SOFT, line=TEAL)
    add_textbox(s, Inches(1.5), Inches(5.35), Inches(4.5), Inches(0.35), "Kayra Çatalkaya", size=15, bold=True, color=NAVY2, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.5), Inches(5.75), Inches(4.5), Inches(0.3), "Kurucu / Ürün", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.3), Inches(5.35), Inches(4.5), Inches(0.35), "AYS Tech", size=15, bold=True, color=NAVY2, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.3), Inches(5.75), Inches(4.5), Inches(0.3), "Altyapı & Teknoloji Partneri", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    slides_meta.append(s)

    total = len(slides_meta)
    for i, slide in enumerate(slides_meta):
        if i == 0:
            continue
        footer_bar(slide, prs, i, total - 1)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    root_copy = ROOT / "KampusteyimAPP_Kurumsal_Tanitim.pptx"
    root_copy.write_bytes(OUT.read_bytes())
    return OUT


if __name__ == "__main__":
    path = build()
    print(path)
    print("bytes", path.stat().st_size)
